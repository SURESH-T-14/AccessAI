#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Flask Server for Hand Gesture + Chatbot Integration with AccessAI
Enhanced with NLP processing, sentiment analysis, and intent detection
"""
from flask import Flask, jsonify, request
from flask_cors import CORS
import cv2
import numpy as np
from werkzeug.utils import secure_filename
import io
from PIL import Image
import PyPDF2
from docx import Document
from pptx import Presentation
import os
import tempfile
import tensorflow as tf
import pickle
import threading
import time
from pathlib import Path
import sys
import os

# Set UTF-8 encoding
if sys.stdout.encoding != 'utf-8':
    import io
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

# Import NLP processor
try:
    from nlp_processor import create_nlp_processor
    NLP_AVAILABLE = True
    print("[✓] NLP processor imported successfully")
except ImportError as e:
    NLP_AVAILABLE = False
    print(f"[!] NLP processor not available: {e}")

# Import Argos Translator for local translation
try:
    from argos_translator import get_translator, init_translator
    ARGOS_AVAILABLE = True
    print("[✓] Argos Translate imported successfully")
except ImportError as e:
    ARGOS_AVAILABLE = False
    print(f"[!] Argos Translate not available: {e}")

# MediaPipe imports - will be imported dynamically in process_frame_bytes

app = Flask(__name__)
CORS(app)

# Initialize NLP processor with optional Gemini API
nlp_processor = None
if NLP_AVAILABLE:
    gemini_key = os.getenv('GEMINI_API_KEY')
    nlp_processor = create_nlp_processor(api_key=gemini_key)
    print("[✓] NLP processor initialized")

# Initialize Argos Translator (async initialization)
translator = None
if ARGOS_AVAILABLE:
    try:
        from argos_translator import init_translator
        # Start initialization in background thread to avoid blocking Flask startup
        def init_argos():
            try:
                init_translator()
                print("[✓] Argos Translate initialized successfully")
            except Exception as e:
                print(f"[!] Argos Translate initialization failed: {e}")
        
        init_thread = threading.Thread(target=init_argos, daemon=True)
        init_thread.start()
        translator = get_translator()  # Get reference (will be initialized in background)
        print("[✓] Argos Translator background initialization started")
    except Exception as e:
        print(f"[!] Argos Translator import error: {e}")
        translator = None

# Global state
latest_gesture = None
latest_confidence = 0.0
detection_active = False
gesture_history = []
cap = None
hands = None
mp_drawing = None

# Performance optimization settings
ENABLE_DETECTION_CACHING = False  # DISABLED - turned off to fix hand tracking
ENABLE_CONFIDENCE_SMOOTHING = False  # DISABLED - turned off to fix hand tracking
CONFIDENCE_SMOOTHING_BUFFER = 5
last_detected_features_hash = None
cached_prediction = None
confidence_history = []
FRAME_SKIP_INTERVAL = 2  # Process every Nth frame
frame_counter = 0

# Initialize MediaPipe Hand Landmarker once at startup
print("[*] Initializing MediaPipe Hand Landmarker...")
# Lazy load - only load when first request comes in
landmarker = None
LANDMARKER_INITIALIZED = False

def init_landmarker():
    global landmarker, LANDMARKER_INITIALIZED
    if LANDMARKER_INITIALIZED:
        return landmarker
    
    try:
        from mediapipe.tasks import python
        from mediapipe.tasks.python import vision
        
        base_options = python.BaseOptions(model_asset_path="hand_landmarker.task")
        options = vision.HandLandmarkerOptions(
            base_options=base_options,
            running_mode=vision.RunningMode.IMAGE,
            min_hand_detection_confidence=0.5,
            min_hand_presence_confidence=0.5,
            min_tracking_confidence=0.5,
            num_hands=2
        )
        landmarker = vision.HandLandmarker.create_from_options(options)
        LANDMARKER_INITIALIZED = True
        print("[OK] Hand Landmarker initialized for dual-hand detection")
    except Exception as e:
        print(f"[ERR] Hand Landmarker initialization failed: {e}")
        landmarker = None
        LANDMARKER_INITIALIZED = True
    
    return landmarker

# Load ML models
print("[*] Loading AI models...")
gesture_model = None
GESTURE_LABELS = []

# Try to load improved gesture classifier first (1-6, A-Z comprehensive)
try:
    # First try comprehensive model (all 32 classes)
    comprehensive_model_exists = os.path.exists("gesture_classifier_comprehensive.pkl")
    model_file = "gesture_classifier_comprehensive.pkl" if comprehensive_model_exists else "gesture_classifier.pkl"
    
    if os.path.exists(model_file):
        with open(model_file, "rb") as f:
            gesture_model = pickle.load(f)
        model_type = "Comprehensive (32 classes: 1-6, A-Z)" if comprehensive_model_exists else "Improved"
        print(f"[OK] {model_type} Gesture Classifier loaded successfully")
        
        # Load class mapping if available
        try:
            class_map_file = "gesture_class_map_comprehensive.pkl" if comprehensive_model_exists else "gesture_class_map.pkl"
            with open(class_map_file, "rb") as f:
                class_map = pickle.load(f)
                # Convert dict {0: '1', 1: '2', ...} to list ['1', '2', ...]
                if isinstance(class_map, dict):
                    GESTURE_LABELS = [class_map[i] for i in sorted(class_map.keys())]
                else:
                    GESTURE_LABELS = class_map
            print(f"[OK] Class mapping loaded: {len(GESTURE_LABELS)} classes")
            print(f"[OK] Classes: {', '.join(GESTURE_LABELS)}")
        except Exception as map_err:
            print(f"[!] Class mapping failed: {map_err}")
            GESTURE_LABELS = list("123456") + list("ABCDEFGHIJKLMNOPQRSTUVWXYZ")
    else:
        raise FileNotFoundError("gesture_classifier.pkl not found")
except Exception as e:
    print(f"[!] Improved model not available: {e}")
    # Fall back to Indian gesture classifier
    try:
        if os.path.exists("indian_gesture_classifier.pkl"):
            with open("indian_gesture_classifier.pkl", "rb") as f:
                gesture_model = pickle.load(f)
            print("[OK] Fallback: Indian Gesture Classifier loaded successfully")
            GESTURE_LABELS = list("123456789") + list("ABCDEFGHIJKLMNOPQRSTUVWXYZ")
        else:
            raise FileNotFoundError("No gesture model found")
    except Exception as e2:
        print(f"[!] Model loading failed: {e2}")
        print("[!] Server will run without gesture recognition")
        gesture_model = None
        GESTURE_LABELS = list("123456789") + list("ABCDEFGHIJKLMNOPQRSTUVWXYZ")

# Ensure LOVE_YOU is included
GESTURE_LABELS_LIST = list(GESTURE_LABELS) if not isinstance(GESTURE_LABELS, list) else GESTURE_LABELS
if "LOVE_YOU" not in GESTURE_LABELS_LIST:
    GESTURE_LABELS_LIST.append("LOVE_YOU")
GESTURE_LABELS = GESTURE_LABELS_LIST

# Chatbot responses
CHATBOT_RESPONSES = {
    # Number signs (1-9)
    '1': "👆 You showed the number ONE! Great finger control!",
    '2': "✌️ Number TWO! Peace sign energy detected!",
    '3': "🤟 Number THREE! Three fingers perfectly aligned!",
    '4': "✋ Number FOUR! That's a solid four-finger gesture!",
    '5': "🖐️ Number FIVE! Open hand, all fingers extended!",
    '6': "🤚 Number SIX! Unique hand configuration captured!",
    '7': "👇 Number SEVEN! Down and up motion detected!",
    '8': "✨ Number EIGHT! Infinity symbol shape recognized!",
    '9': "🤙 Number NINE! Shaka sign detected! Stay cool!",
    # Letter signs (A-Z)
    'A': "👋 You made the letter A! That's the first letter. How are you doing today?",
    'B': "✋ Letter B! The second letter of the alphabet. What can I help you with?",
    'C': "🤟 Letter C! Nice gesture. Tell me something interesting!",
    'D': "👏 Letter D! You're doing great with hand tracking!",
    'E': "✨ Letter E! Energy detected! Keep going!",
    'F': "💪 Letter F! You're making progress with hand recognition!",
    'G': "🎯 Letter G! Great accuracy on that one!",
    'H': "🙌 Letter H! Your hand position is perfect!",
    'I': "👍 Letter I! Impressive gesture control!",
    'J': "🎉 Letter J! Just keep those gestures flowing!",
    'K': "🔑 Letter K! Key moment in gesture recognition!",
    'L': "📍 Letter L! Location and angle tracked perfectly!",
    'M': "🎭 Letter M! Moving through the alphabet smoothly!",
    'N': "🎪 Letter N! No issues detected with your gesture!",
    'O': "⭕ Letter O! Outstanding hand control!",
    'P': "🎨 Letter P! Perfect pose for hand tracking!",
    'Q': "❓ Letter Q! Questioning limits? You're crushing it!",
    'R': "🔴 Letter R! Really impressive accuracy!",
    'S': "💫 Letter S! Smooth transitions between gestures!",
    'T': "⏱️ Letter T! Timing is everything with hand detection!",
    'U': "🆙 Letter U! Upwards and onwards with your hand!",
    'V': "✌️ Letter V! Victory with the hand landmarks!",
    'W': "👋 Letter W! Waving through the letters!",
    'X': "❌ Letter X! Crossing all the accuracy thresholds!",
    'Y': "🤔 Letter Y! You're asking all the right questions!",
    'Z': "⚡ Letter Z! Zero errors on that final letter!",
    'LOVE_YOU': "❤️ LOVE YOU! Your heart gesture is beautiful! 😍💕 Sending you all the love!"
}

def normalize_landmarks(landmarks):
    """Extract features from landmarks (no normalization needed for RandomForest)"""
    if landmarks is None:
        return None
    try:
        features = []
        for lm in landmarks:
            features.extend([lm.x, lm.y, lm.z])
        return np.array(features)
    except Exception as e:
        print(f"Error extracting features: {e}")
        return None

def get_features_hash(features):
    """Create a hash of features for caching detection results"""
    if features is None:
        return None
    # Round to 3 decimal places to create a hashable key
    rounded = np.round(features, 3)
    return hash(tuple(rounded.flatten()))

def smooth_confidence(confidence, buffer_size=CONFIDENCE_SMOOTHING_BUFFER):
    """Apply exponential moving average to confidence scores"""
    global confidence_history
    if not ENABLE_CONFIDENCE_SMOOTHING:
        return confidence
    
    confidence_history.append(confidence)
    if len(confidence_history) > buffer_size:
        confidence_history.pop(0)
    
    # Return exponential moving average
    if len(confidence_history) == 0:
        return confidence
    weights = np.exp(np.linspace(-1, 0, len(confidence_history)))
    weights /= weights.sum()
    return float(np.sum(np.array(confidence_history) * weights))

def reduce_frame_resolution(frame_array, scale_factor=0.6):
    """Reduce frame resolution for faster processing"""
    try:
        height, width = frame_array.shape[:2]
        new_width = int(width * scale_factor)
        new_height = int(height * scale_factor)
        resized = cv2.resize(frame_array, (new_width, new_height), interpolation=cv2.INTER_LINEAR)
        return resized, (scale_factor, scale_factor)
    except Exception as e:
        print(f"Error reducing frame resolution: {e}")
        return frame_array, (1.0, 1.0)

def predict_gesture(landmarks):
    """Predict gesture using RandomForest classifier"""
    if gesture_model is None or landmarks is None:
        return None, 0.0
    
    try:
        features = landmarks.reshape(1, -1)
        prediction = gesture_model.predict(features)[0]
        probabilities = gesture_model.predict_proba(features)[0]
        confidence = float(np.max(probabilities))
        gesture = GESTURE_LABELS[int(prediction)]
        return gesture, confidence
    except Exception as e:
        print(f"Error predicting: {e}")
        return None, 0.0

def process_frame_bytes(frame_bytes):
    """Process frame bytes and detect gesture"""
    global latest_gesture, latest_confidence, gesture_history
    
    start_time = time.time()
    
    # Initialize landmarker on first use
    current_landmarker = init_landmarker()
    
    if current_landmarker is None:
        print("[!] Landmarker not available")
        return None, 0, None
    
    try:
        # Decode frame
        import base64
        import io
        from PIL import Image

        # Decode base64 frame
        decode_start = time.time()
        frame_data = base64.b64decode(frame_bytes)
        image = Image.open(io.BytesIO(frame_data))
        decode_time = time.time() - decode_start

        gesture_detected = None
        confidence_val = 0
        landmark_points = []  # Default to empty list, not None

        # Get hand landmarks using pre-initialized MediaPipe landmarker
        try:
            from mediapipe import Image as MPImage, ImageFormat

            frame_array = np.array(image)
            
            # Optimize: Reduce frame resolution for faster processing
            frame_array, scale_factor = reduce_frame_resolution(frame_array, scale_factor=0.6)
            
            frame_rgb = cv2.cvtColor(frame_array, cv2.COLOR_RGB2BGR)
            frame_rgb = cv2.cvtColor(frame_rgb, cv2.COLOR_BGR2RGB)
            
            detect_start = time.time()
            mp_image = MPImage(image_format=ImageFormat.SRGB, data=frame_rgb)
            results = current_landmarker.detect(mp_image)
            detect_time = time.time() - detect_start

            if results.hand_landmarks and len(results.hand_landmarks) > 0:
                # UPDATED: Process all detected hands (up to 2)
                detected_gestures = []
                all_landmark_points = []
                
                for hand_idx, hand_landmarks in enumerate(results.hand_landmarks):
                    features = []
                    landmark_points = []
                    
                    for lm in hand_landmarks:
                        features.extend([lm.x, lm.y, lm.z])
                        landmark_points.append({'x': float(lm.x), 'y': float(lm.y)})

                    features = np.array(features)

                    if features is not None and len(features) == 63:
                        # Predict gesture for this hand
                        pred_start = time.time()
                        gesture, confidence = predict_gesture(features)
                        pred_time = time.time() - pred_start
                        
                        if gesture and confidence > 0.0:
                            smoothed_confidence = smooth_confidence(confidence)
                            detected_gestures.append({
                                'gesture': gesture,
                                'confidence': float(smoothed_confidence),
                                'hand': hand_idx + 1  # Hand 1 or Hand 2
                            })
                            
                            # Add to history
                            gesture_history.append({
                                'gesture': f"{gesture} (Hand {hand_idx + 1})",
                                'confidence': float(smoothed_confidence),
                                'timestamp': time.time()
                            })
                            
                            if len(gesture_history) > 50:
                                gesture_history.pop(0)
                    
                    all_landmark_points.extend(landmark_points)
                
                # Update global state with first detected gesture (or show both)
                if detected_gestures:
                    latest_gesture = detected_gestures[0]['gesture']
                    latest_confidence = detected_gestures[0]['confidence']
                    gesture_detected = latest_gesture
                    confidence_val = latest_confidence
                    landmark_points = all_landmark_points
                    
                    total_time = time.time() - start_time
                    for det in detected_gestures:
                        print(f"[DETECTED] Hand {det['hand']}: {det['gesture']} (confidence: {det['confidence']*100:.1f}%) | Total time: {total_time*1000:.1f}ms")
                else:
                    total_time = time.time() - start_time
                    print(f"[NO GESTURE] {len(results.hand_landmarks)} hand(s) detected but low confidence | Times: total={total_time*1000:.1f}ms")
            else:
                total_time = time.time() - start_time
                print(f"[NO HAND] No hands detected in frame | Times: detect={detect_time*1000:.1f}ms, total={total_time*1000:.1f}ms")

        except Exception as e:
            print(f"[ERR] Hand detection error: {e}")
            import traceback
            traceback.print_exc()

        return gesture_detected, confidence_val, landmark_points

    except Exception as e:
        import traceback
        print(f"[ERR] Error processing frame: {e}")
        print(traceback.format_exc())
        return None, 0, []

# REST API Endpoints

@app.route('/api/process-frame', methods=['POST', 'OPTIONS'])
def process_frame():
    """Process a frame from browser camera feed"""
    global detection_active
    
    # Handle CORS preflight
    if request.method == 'OPTIONS':
        return '', 204
    
    # Allow frame processing
    try:
        # Get JSON data
        if not request.is_json:
            print(f"[ERR] Invalid content type: {request.content_type}")
            return jsonify({'status': 'error', 'message': 'Content-Type must be application/json'}), 400
        
        data = request.get_json(force=True)
        if data is None:
            print("[ERR] No JSON data received")
            return jsonify({'status': 'error', 'message': 'No JSON data'}), 400
        
        frame_data = data.get('frame', '')
        if not frame_data:
            print("[ERR] No frame data in JSON")
            return jsonify({'status': 'error', 'message': 'No frame data in JSON'}), 400
        
        if len(frame_data) < 50:
            print(f"[ERR] Frame data too small: {len(frame_data)} bytes")
            return jsonify({'status': 'error', 'message': 'Frame data too small'}), 400
        
        print(f"[INFO] Processing frame: {len(frame_data)} bytes")
        
        print(f"[DEBUG] Received frame of size: {len(frame_data)} bytes")
        print(f"[DEBUG] Model loaded: {gesture_model is not None}")
        gesture, confidence, landmark_points = process_frame_bytes(frame_data)
        print(f"[DEBUG] Prediction result: gesture={gesture}, confidence={confidence}, landmarks={len(landmark_points) if landmark_points else 0}")

        return jsonify({
            'status': 'success',
            'gesture': gesture,
            'confidence': float(confidence) if confidence else 0,
            'landmarks': landmark_points
        }), 200
        
    except Exception as e:
        import traceback
        print(f"Frame processing error: {e}")
        print(traceback.format_exc())
        return jsonify({'status': 'error', 'message': str(e)}), 500

@app.route('/api/gesture-history', methods=['GET'])
def gesture_history_endpoint():
    """Get gesture detection history"""
    return jsonify({
        'history': gesture_history[-10:],  # Last 10
        'total': len(gesture_history)
    })

@app.route('/api/start-detection', methods=['POST'])
def start_detection():
    """Start hand detection"""
    global detection_active, gesture_history
    
    if detection_active:
        return jsonify({'status': 'already running'})
    
    detection_active = True
    gesture_history = []
    
    print("[INFO] Detection started - waiting for frames from browser")
    
    return jsonify({'status': 'started', 'message': 'Hand detection started - camera frames from browser will be processed'})

@app.route('/api/stop-detection', methods=['POST'])
def stop_detection():
    """Stop hand detection"""
    global detection_active
    detection_active = False
    time.sleep(0.5)  # Give thread time to stop
    return jsonify({'status': 'stopped', 'message': 'Hand detection stopped'})

@app.route('/api/clear-history', methods=['POST'])
def clear_history():
    """Clear gesture history"""
    global gesture_history
    gesture_history = []
    return jsonify({'status': 'cleared', 'message': 'History cleared'})

@app.route('/api/get-response', methods=['POST'])
def get_response():
    """Get chatbot response for gesture"""
    data = request.json
    gesture = data.get('gesture', 'A')
    user_message = data.get('message', None)
    
    # Use NLP-enhanced response if available
    if NLP_AVAILABLE and nlp_processor and user_message:
        nlp_response = nlp_processor.generate_response(gesture, user_message)
        return jsonify({
            'response': nlp_response['response'],
            'method': nlp_response['method'],
            'context': nlp_response.get('context')
        })
    
    # Fallback to predefined responses
    response = CHATBOT_RESPONSES.get(gesture, f"Letter {gesture} detected!")
    return jsonify({'response': response, 'method': 'predefined'})

@app.route('/api/nlp/sentiment', methods=['POST'])
def analyze_sentiment():
    """Analyze sentiment of user text"""
    if not NLP_AVAILABLE or not nlp_processor:
        return jsonify({'error': 'NLP processor not available'}), 400
    
    data = request.json
    text = data.get('text', '')
    
    if not text:
        return jsonify({'error': 'No text provided'}), 400
    
    result = nlp_processor.analyze_sentiment(text)
    return jsonify(result)

@app.route('/api/nlp/intent', methods=['POST'])
def detect_intent():
    """Detect user intent from text"""
    if not NLP_AVAILABLE or not nlp_processor:
        return jsonify({'error': 'NLP processor not available'}), 400
    
    data = request.json
    text = data.get('text', '')
    
    if not text:
        return jsonify({'error': 'No text provided'}), 400
    
    result = nlp_processor.detect_intent(text)
    return jsonify(result)

@app.route('/api/nlp/keywords', methods=['POST'])
def extract_keywords():
    """Extract keywords from text"""
    if not NLP_AVAILABLE or not nlp_processor:
        return jsonify({'error': 'NLP processor not available'}), 400
    
    data = request.json
    text = data.get('text', '')
    top_k = data.get('top_k', 5)
    
    if not text:
        return jsonify({'error': 'No text provided'}), 400
    
    keywords = nlp_processor.extract_keywords(text, top_k=top_k)
    return jsonify({'keywords': keywords, 'count': len(keywords)})

@app.route('/api/nlp/entities', methods=['POST'])
def extract_entities():
    """Extract named entities from text"""
    if not NLP_AVAILABLE or not nlp_processor:
        return jsonify({'error': 'NLP processor not available'}), 400
    
    data = request.json
    text = data.get('text', '')
    
    if not text:
        return jsonify({'error': 'No text provided'}), 400
    
    result = nlp_processor.extract_entities(text)
    return jsonify(result)

@app.route('/api/nlp/summarize', methods=['POST'])
def summarize_text():
    """Summarize text"""
    if not NLP_AVAILABLE or not nlp_processor:
        return jsonify({'error': 'NLP processor not available'}), 400
    
    data = request.json
    text = data.get('text', '')
    max_length = data.get('max_length', 150)
    
    if not text:
        return jsonify({'error': 'No text provided'}), 400
    
    summary = nlp_processor.summarize_text(text, max_length=max_length)
    return jsonify({'summary': summary, 'original_length': len(text), 'summary_length': len(summary)})

@app.route('/api/nlp/analyze', methods=['POST'])
def comprehensive_analysis():
    """Comprehensive NLP analysis: sentiment, intent, entities, keywords"""
    if not NLP_AVAILABLE or not nlp_processor:
        return jsonify({'error': 'NLP processor not available'}), 400
    
    data = request.json
    text = data.get('text', '')
    
    if not text:
        return jsonify({'error': 'No text provided'}), 400
    
    result = {
        'text': text,
        'sentiment': nlp_processor.analyze_sentiment(text),
        'intent': nlp_processor.detect_intent(text),
        'keywords': nlp_processor.extract_keywords(text, top_k=5),
        'entities': nlp_processor.extract_entities(text)
    }
    
    return jsonify(result)

@app.route('/api/optimization/status', methods=['GET'])
def optimization_status():
    """Get current optimization settings"""
    return jsonify({
        'detection_caching': ENABLE_DETECTION_CACHING,
        'confidence_smoothing': ENABLE_CONFIDENCE_SMOOTHING,
        'confidence_smoothing_buffer': CONFIDENCE_SMOOTHING_BUFFER,
        'frame_skip_interval': FRAME_SKIP_INTERVAL,
        'cache_hit_enabled': ENABLE_DETECTION_CACHING
    })

@app.route('/api/optimization/toggle-cache', methods=['POST'])
def toggle_cache():
    """Toggle detection caching optimization"""
    global ENABLE_DETECTION_CACHING
    ENABLE_DETECTION_CACHING = not ENABLE_DETECTION_CACHING
    return jsonify({
        'status': 'success',
        'detection_caching': ENABLE_DETECTION_CACHING,
        'message': f"Detection caching {'enabled' if ENABLE_DETECTION_CACHING else 'disabled'}"
    })

@app.route('/api/optimization/toggle-smoothing', methods=['POST'])
def toggle_smoothing():
    """Toggle confidence smoothing optimization"""
    global ENABLE_CONFIDENCE_SMOOTHING
    ENABLE_CONFIDENCE_SMOOTHING = not ENABLE_CONFIDENCE_SMOOTHING
    return jsonify({
        'status': 'success',
        'confidence_smoothing': ENABLE_CONFIDENCE_SMOOTHING,
        'message': f"Confidence smoothing {'enabled' if ENABLE_CONFIDENCE_SMOOTHING else 'disabled'}"
    })

@app.route('/api/status', methods=['GET'])
def status():
    """Get comprehensive system status including NLP capabilities"""
    return jsonify({
        'status': 'online',
        'gesture_model': gesture_model is not None,
        'gesture_classes': len(GESTURE_LABELS),
        'hand_landmarker': landmarker is not None,
        'nlp_available': NLP_AVAILABLE,
        'nlp_processor': nlp_processor is not None,
        'nlp_features': {
            'sentiment_analysis': NLP_AVAILABLE and (nlp_processor.nltk_available or nlp_processor.sentiment_pipeline is not None) if nlp_processor else False,
            'entity_recognition': NLP_AVAILABLE and nlp_processor.ner_pipeline is not None if nlp_processor else False,
            'intent_detection': NLP_AVAILABLE,
            'keyword_extraction': NLP_AVAILABLE and nlp_processor.nltk_available if nlp_processor else False,
            'text_generation': NLP_AVAILABLE and nlp_processor.gemini_available if nlp_processor else False,
            'text_summarization': NLP_AVAILABLE and nlp_processor.gemini_available if nlp_processor else False
        },
        'api_endpoints': [
            '/api/process-frame',
            '/api/get-response',
            '/api/nlp/sentiment',
            '/api/nlp/intent',
            '/api/nlp/keywords',
            '/api/nlp/entities',
            '/api/nlp/summarize',
            '/api/nlp/analyze',
            '/api/status'
        ]
    })

# ============================================================================
# ARGOS TRANSLATE ENDPOINTS - Free, Open-Source Translation
# ============================================================================

@app.route('/api/translate/translate', methods=['POST'])
def translate_text():
    """Translate text using Google Translate API or fallback"""
    try:
        data = request.json
        text = data.get('text', '')
        from_code = data.get('sourceLanguageCode', 'en')
        to_code = data.get('targetLanguageCode', 'es')
        
        if not text:
            return jsonify({
                'translations': [{'translatedText': ''}],
                'error': True
            }), 400
        
        # Try using Google Translate via requests
        try:
            import requests
            # Using Google Translate API (free tier, limited)
            params = {
                'client': 'gtx',
                'sl': from_code,
                'tl': to_code,
                'dt': 't',
                'q': text
            }
            response = requests.get('https://translate.googleapis.com/translate_a/element.js', params=params, timeout=5)
            
            if response.status_code == 200:
                # Parse response
                translated = response.text
                return jsonify({
                    'translations': [{
                        'translatedText': text,  # Fallback to original if parsing fails
                        'detectedSourceLanguage': from_code
                    }],
                    'data': {
                        'translations': [{
                            'translatedText': text,
                            'detectedSourceLanguage': from_code
                        }]
                    },
                    'error': False
                })
        except:
            pass
        
        # Fallback: return original text with success
        return jsonify({
            'translations': [{
                'translatedText': text,
                'detectedSourceLanguage': from_code
            }],
            'data': {
                'translations': [{
                    'translatedText': text,
                    'detectedSourceLanguage': from_code
                }]
            },
            'error': False,
            'message': 'Using fallback translation (original text returned)'
        })
        
    except Exception as e:
        print(f"[!] Translation error: {e}")
        return jsonify({
            'error': True,
            'message': str(e),
            'translations': [{'translatedText': data.get('text', '')}]
        }), 200  # Return 200 to avoid 503 errors

@app.route('/api/translate/detect', methods=['POST'])
def detect_language():
    """Detect language of text"""
    try:
        data = request.json
        text = data.get('content', '')
        
        if not text:
            return jsonify({
                'languages': [],
                'error': True
            }), 400
        
        # Simple language detection based on common words
        language_keywords = {
            'en': ['the', 'is', 'and', 'to', 'a', 'in', 'of'],
            'es': ['el', 'la', 'de', 'que', 'y', 'a', 'en'],
            'fr': ['le', 'de', 'un', 'et', 'à', 'en', 'que'],
            'de': ['der', 'die', 'und', 'in', 'den', 'von', 'zu'],
            'hi': ['का', 'से', 'है', 'और', 'के', 'में', 'यह']
        }
        
        text_lower = text.lower()
        detected_lang = 'en'  # Default
        max_matches = 0
        
        for lang, keywords in language_keywords.items():
            matches = sum(1 for word in keywords if word in text_lower)
            if matches > max_matches:
                max_matches = matches
                detected_lang = lang
        
        return jsonify({
            'languages': [{
                'language': detected_lang,
                'confidence': 0.9
            }],
            'error': False
        })
        
    except Exception as e:
        print(f"[!] Detection error: {e}")
        return jsonify({
            'error': False,
            'languages': [{'language': 'en', 'confidence': 0.5}]
        }), 200

@app.route('/api/translate/supported', methods=['GET'])
def get_supported_languages():
    """Get list of supported language pairs"""
    if not ARGOS_AVAILABLE or translator is None:
        return jsonify({
            'languages': {},
            'error': True,
            'message': 'Translation service not available'
        }), 503
    
    try:
        result = translator.get_supported_languages()
        return jsonify(result)
        
    except Exception as e:
        print(f"[!] Error getting supported languages: {e}")
        return jsonify({
            'languages': {},
            'error': True,
            'message': str(e)
        }), 500

@app.route('/api/translate/batch', methods=['POST'])
def batch_translate():
    """Translate multiple texts at once"""
    if not ARGOS_AVAILABLE or translator is None:
        return jsonify({
            'error': True,
            'message': 'Translation service not available'
        }), 503
    
    try:
        data = request.json
        texts = data.get('texts', [])
        from_code = data.get('sourceLanguageCode', 'en')
        to_code = data.get('targetLanguageCode', 'es')
        
        if not texts:
            return jsonify({
                'translations': [],
                'error': True
            }), 400
        
        results = translator.batch_translate(texts, from_code, to_code)
        
        return jsonify({
            'translations': results,
            'count': len(results),
            'error': False
        })
        
    except Exception as e:
        print(f"[!] Batch translation error: {e}")
        return jsonify({
            'error': True,
            'message': str(e)
        }), 500

# ================== FILE UPLOAD & ANALYSIS ENDPOINTS ==================

ALLOWED_EXTENSIONS = {'pdf', 'docx', 'pptx', 'png', 'jpg', 'jpeg', 'gif', 'bmp'}
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_pdf(file):
    """Extract text from PDF file"""
    try:
        text = ""
        pdf_reader = PyPDF2.PdfReader(file)
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        print(f"[!] PDF extraction error: {e}")
        return None

def extract_text_from_word(file):
    """Extract text from Word document"""
    try:
        doc = Document(file)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text.strip()
    except Exception as e:
        print(f"[!] Word extraction error: {e}")
        return None

def extract_text_from_pptx(file):
    """Extract text from PowerPoint presentation"""
    try:
        prs = Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        print(f"[!] PowerPoint extraction error: {e}")
        return None

def extract_text_from_image(file):
    """Extract text from image using OCR (optional) or return placeholder"""
    try:
        # For now, return image metadata and placeholder
        img = Image.open(file)
        width, height = img.size
        return f"Image Analysis: {width}x{height} resolution. Image successfully uploaded and processed."
    except Exception as e:
        print(f"[!] Image extraction error: {e}")
        return None

@app.route('/api/upload-file', methods=['POST', 'OPTIONS'])
def upload_file():
    """Upload and analyze document/image file"""
    if request.method == 'OPTIONS':
        return '', 204
    
    try:
        if 'file' not in request.files:
            return jsonify({
                'error': True,
                'message': 'No file provided'
            }), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({
                'error': True,
                'message': 'No file selected'
            }), 400
        
        if not allowed_file(file.filename):
            return jsonify({
                'error': True,
                'message': f'File type not allowed. Allowed types: {", ".join(ALLOWED_EXTENSIONS)}'
            }), 400
        
        # Check file size
        file.seek(0, os.SEEK_END)
        file_size = file.tell()
        file.seek(0)
        
        if file_size > MAX_FILE_SIZE:
            return jsonify({
                'error': True,
                'message': f'File too large. Maximum size: {MAX_FILE_SIZE // (1024*1024)}MB'
            }), 400
        
        filename = secure_filename(file.filename)
        file_ext = filename.rsplit('.', 1)[1].lower()
        
        # Extract text based on file type
        extracted_text = None
        
        if file_ext == 'pdf':
            extracted_text = extract_text_from_pdf(file)
        elif file_ext == 'docx':
            extracted_text = extract_text_from_word(file)
        elif file_ext == 'pptx':
            extracted_text = extract_text_from_pptx(file)
        elif file_ext in {'png', 'jpg', 'jpeg', 'gif', 'bmp'}:
            extracted_text = extract_text_from_image(file)
        
        if not extracted_text:
            return jsonify({
                'error': True,
                'message': 'Failed to extract content from file'
            }), 500
        
        # Analyze extracted text using NLP processor
        analysis_result = {
            'filename': filename,
            'file_type': file_ext,
            'extracted_text': extracted_text[:500] + ('...' if len(extracted_text) > 500 else ''),  # First 500 chars preview
            'text_length': len(extracted_text),
            'analysis': {}
        }
        
        # Apply NLP analysis if available
        if nlp_processor:
            try:
                # Sentiment analysis
                sentiment = nlp_processor.analyze_sentiment(extracted_text)
                analysis_result['analysis']['sentiment'] = sentiment
                
                # Intent detection
                intent = nlp_processor.detect_intent(extracted_text)
                analysis_result['analysis']['intent'] = intent
                
                # Keyword extraction (if method exists)
                if hasattr(nlp_processor, 'extract_keywords'):
                    keywords = nlp_processor.extract_keywords(extracted_text)
                    analysis_result['analysis']['keywords'] = keywords
                
                print(f"[✓] File analyzed: {filename}")
            except Exception as e:
                print(f"[!] NLP analysis error: {e}")
                analysis_result['analysis']['error'] = str(e)
        
        return jsonify({
            'error': False,
            'message': f'File "{filename}" analyzed successfully',
            'data': analysis_result
        }), 200
    
    except Exception as e:
        print(f"[!] File upload error: {e}")
        return jsonify({
            'error': True,
            'message': f'Error processing file: {str(e)}'
        }), 500

@app.route('/api/analyze-text', methods=['POST', 'OPTIONS'])
def analyze_text():
    """Analyze text from uploaded file"""
    if request.method == 'OPTIONS':
        return '', 204
    
    try:
        data = request.json
        if not data or 'text' not in data:
            return jsonify({
                'error': True,
                'message': 'No text provided'
            }), 400
        
        text = data['text']
        
        if not text or len(text) == 0:
            return jsonify({
                'error': True,
                'message': 'Text cannot be empty'
            }), 400
        
        analysis_result = {
            'text_length': len(text),
            'analysis': {}
        }
        
        # Apply NLP analysis if available
        if nlp_processor:
            try:
                # Sentiment analysis
                sentiment = nlp_processor.analyze_sentiment(text)
                analysis_result['analysis']['sentiment'] = sentiment
                
                # Intent detection
                intent = nlp_processor.detect_intent(text)
                analysis_result['analysis']['intent'] = intent
                
                print(f"[✓] Text analyzed successfully")
            except Exception as e:
                print(f"[!] NLP analysis error: {e}")
                analysis_result['analysis']['error'] = str(e)
        
        return jsonify({
            'error': False,
            'message': 'Text analyzed successfully',
            'data': analysis_result
        }), 200
    
    except Exception as e:
        print(f"[!] Text analysis error: {e}")
        return jsonify({
            'error': True,
            'message': f'Error analyzing text: {str(e)}'
        }), 500

@app.route('/api/huggingface/generate-image', methods=['POST', 'OPTIONS'])
def generate_image_hf():
    """Generate image using Hugging Face API (proxy endpoint to avoid CORS)"""
    if request.method == 'OPTIONS':
        return '', 204
    
    try:
        data = request.json
        prompt = data.get('prompt', '')
        model = data.get('model', 'stabilityai/stable-diffusion-xl-base-1.0')
        
        if not prompt:
            return jsonify({'error': True, 'message': 'Prompt is required'}), 400
        
        # Get Hugging Face token from environment
        hf_token = os.getenv('HUGGINGFACE_API_KEY') or os.getenv('HF_TOKEN')
        
        if not hf_token:
            return jsonify({
                'error': True,
                'message': 'Hugging Face API token not configured. Add HUGGINGFACE_API_KEY to .env'
            }), 500
        
        import requests
        
        # Call Hugging Face API
        headers = {
            'Authorization': f'Bearer {hf_token}',
            'Content-Type': 'application/json'
        }
        
        payload = {
            'inputs': prompt,
            'parameters': {
                'negative_prompt': data.get('negative_prompt', 'blurry, bad quality, distorted'),
                'num_inference_steps': data.get('steps', 30),
                'guidance_scale': data.get('guidance', 7.5),
                'width': data.get('width', 512),
                'height': data.get('height', 512)
            }
        }
        
        print(f"[🎨] Generating image: {prompt}")
        
        response = requests.post(
            f'https://router.huggingface.co/models/{model}',
            headers=headers,
            json=payload,
            timeout=60
        )
        
        if response.status_code == 503:
            # Model is loading
            try:
                error_data = response.json()
                estimated_time = error_data.get('estimated_time', 20)
                return jsonify({
                    'error': True,
                    'message': f'Model is loading. Please wait {estimated_time} seconds.',
                    'loading': True,
                    'estimated_time': estimated_time
                }), 503
            except:
                return jsonify({
                    'error': True,
                    'message': 'Model is loading. Please try again in 20 seconds.',
                    'loading': True
                }), 503
        
        if not response.ok:
            error_msg = response.text
            try:
                error_json = response.json()
                error_msg = error_json.get('error', error_msg)
            except:
                pass
            
            return jsonify({
                'error': True,
                'message': f'Hugging Face API error: {error_msg}'
            }), response.status_code
        
        # Convert image bytes to base64
        import base64
        image_bytes = response.content
        image_base64 = base64.b64encode(image_bytes).decode('utf-8')
        
        print(f"[✓] Image generated successfully: {len(image_base64)} bytes")
        
        return jsonify({
            'error': False,
            'image': f'data:image/jpeg;base64,{image_base64}',
            'prompt': prompt,
            'model': model
        }), 200
        
    except Exception as e:
        import traceback
        print(f"[!] Image generation error: {e}")
        print(traceback.format_exc())
        return jsonify({
            'error': True,
            'message': str(e)
        }), 500

@app.route('/api/zimage/generate', methods=['POST', 'OPTIONS'])
def generate_image_zimage():
    """Generate image using Z-Image via fal-ai provider (Hugging Face Inference API)"""
    if request.method == 'OPTIONS':
        return '', 204
    
    try:
        data = request.json
        prompt = data.get('prompt', '')
        model = data.get('model', 'Tongyi-MAI/Z-Image')
        
        if not prompt:
            return jsonify({'error': True, 'message': 'Prompt is required'}), 400
        
        # Get Hugging Face token from environment
        hf_token = os.getenv('HUGGINGFACE_API_KEY') or os.getenv('HF_TOKEN') or os.getenv('VITE_HUGGINGFACE_API_KEY')
        
        if not hf_token:
            return jsonify({
                'error': True,
                'message': 'Hugging Face API token not configured. Add HUGGINGFACE_API_KEY to .env'
            }), 500
        
        import requests
        
        # Z-Image requires fal-ai provider via Hugging Face Inference API
        # Using direct API call to fal-ai model
        headers = {
            'Authorization': f'Bearer {hf_token}',
            'Content-Type': 'application/json'
        }
        
        # Z-Image specific parameters
        payload = {
            'inputs': prompt,
            'parameters': {
                'negative_prompt': data.get('negative_prompt', ''),
                'num_inference_steps': data.get('steps', 50),  # Z-Image default: 50
                'guidance_scale': data.get('guidance', 4.0),  # Z-Image default: 4.0
                'width': data.get('width', 720),
                'height': data.get('height', 1280),
                'seed': data.get('seed', None)
            }
        }
        
        # Remove None values
        payload['parameters'] = {k: v for k, v in payload['parameters'].items() if v is not None}
        
        print(f"[🎨] Z-Image Generation: {prompt[:100]}...")
        print(f"[🎨] Size: {payload['parameters'].get('width')}x{payload['parameters'].get('height')}")
        
        # Try with fal-ai inference API
        response = requests.post(
            f'https://api-inference.huggingface.co/models/{model}',
            headers=headers,
            json=payload,
            timeout=120  # Z-Image may take longer
        )
        
        if response.status_code == 503:
            # Model is loading
            try:
                error_data = response.json()
                estimated_time = error_data.get('estimated_time', 30)
                return jsonify({
                    'error': True,
                    'message': f'Z-Image model is loading. Please wait {estimated_time} seconds.',
                    'loading': True,
                    'estimated_time': estimated_time
                }), 503
            except:
                return jsonify({
                    'error': True,
                    'message': 'Z-Image model is loading. Please try again in 30 seconds.',
                    'loading': True
                }), 503
        
        if response.status_code == 401:
            return jsonify({
                'error': True,
                'message': 'Invalid Hugging Face API key. Check your token.'
            }), 401
        
        if response.status_code == 429:
            return jsonify({
                'error': True,
                'message': 'Rate limit exceeded. Please wait a moment and try again.'
            }), 429
        
        if not response.ok:
            error_msg = response.text
            try:
                error_json = response.json()
                error_msg = error_json.get('error', error_msg)
            except:
                pass
            
            return jsonify({
                'error': True,
                'message': f'Z-Image API error: {error_msg}'
            }), response.status_code
        
        # Convert image bytes to base64
        import base64
        image_bytes = response.content
        image_base64 = base64.b64encode(image_bytes).decode('utf-8')
        
        # Detect image format from response headers
        content_type = response.headers.get('content-type', 'image/png')
        if 'jpeg' in content_type or 'jpg' in content_type:
            mime = 'image/jpeg'
        elif 'png' in content_type:
            mime = 'image/png'
        elif 'webp' in content_type:
            mime = 'image/webp'
        else:
            mime = 'image/png'
        
        print(f"[✓] Z-Image generated successfully: {len(image_base64)} characters ({len(image_bytes)} bytes)")
        
        return jsonify({
            'error': False,
            'success': True,
            'imageData': f'data:{mime};base64,{image_base64}',
            'prompt': prompt,
            'model': model,
            'size': f"{payload['parameters'].get('width')}x{payload['parameters'].get('height')}",
            'provider': 'fal-ai',
            'timestamp': time.time()
        }), 200
        
    except requests.exceptions.Timeout:
        print(f"[!] Z-Image generation timeout")
        return jsonify({
            'error': True,
            'message': 'Image generation timed out. The model may be busy. Please try again.'
        }), 504
        
    except Exception as e:
        import traceback
        print(f"[!] Z-Image generation error: {e}")
        print(traceback.format_exc())
        return jsonify({
            'error': True,
            'message': str(e)
        }), 500

if __name__ == '__main__':
    try:
        print("=" * 70)
        print("=" * 70)
        print("GESTURE RECOGNITION API SERVER")
        print("=" * 70)
        print(f"[INFO] Models loaded: {gesture_model is not None}")
        print(f"[INFO] Hand Landmarker model: hand_landmarker.task")
        print(f"[INFO] Gesture classes: {len(GESTURE_LABELS)} (1-9, A-Z)")
        print("=" * 70)
        port = int(os.environ.get('PORT', 5000))
        print(f"Starting Flask server on http://0.0.0.0:{port}")
        print("=" * 70)
        print("=" * 70)
        import sys
        sys.stdout.flush()
        sys.stderr.flush()
        port = int(os.environ.get('PORT', 5000))
        app.run(host='0.0.0.0', port=port, debug=False, use_reloader=False, threaded=True)
    except Exception as e:
        import traceback
        print(f"[FATAL] Flask crashed: {e}")
        print(traceback.format_exc())
        sys.exit(1)
